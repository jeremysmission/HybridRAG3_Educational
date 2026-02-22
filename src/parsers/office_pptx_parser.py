# ============================================================================
# HybridRAG -- PPTX Parser (src/parsers/office_pptx_parser.py)
# ============================================================================
#
# WHAT THIS FILE DOES (plain English):
#   Reads PowerPoint (.pptx) presentation files and extracts all the text
#   from every slide. Each text block is tagged with its slide number
#   so you can tell where the information came from when searching.
#
# WHY THIS MATTERS:
#   Briefings, training materials, and technical presentations often
#   contain critical information that exists nowhere else. Making
#   slide content searchable means you can find "What frequency range
#   was mentioned in the system briefing?" without opening 50 decks.
#
# HOW IT WORKS:
#   1. Open the .pptx file using python-pptx
#   2. Loop through every slide in order
#   3. For each slide, loop through every "shape" (text boxes, titles,
#      bullets, tables that have text)
#   4. Tag each text block with [SLIDE N] so the source is traceable
#   5. Return combined text + details (slide count, text block count)
#
# INTERNET ACCESS: NONE
# ============================================================================

from __future__ import annotations

from pathlib import Path
from typing import Tuple, Dict, Any


class PptxParser:
    def parse(self, file_path: str) -> str:
        text, _ = self.parse_with_details(file_path)
        return text

    def parse_with_details(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        path = Path(file_path)
        details: Dict[str, Any] = {"file": str(path), "parser": "PptxParser"}

        try:
            from pptx import Presentation  # python-pptx
        except Exception as e:
            details["error"] = f"IMPORT_ERROR: {type(e).__name__}: {e}"
            return "", details

        try:
            pres = Presentation(str(path))

            parts = []
            slide_text_count = 0

            # Walk through every slide in presentation order.
            # Each slide can contain many "shapes" -- text boxes, titles,
            # bullet lists, table cells, etc. We grab all of them.
            for si, slide in enumerate(pres.slides):
                for shape in slide.shapes:
                    # Some shapes are images or charts with no text property
                    if not hasattr(shape, "text"):
                        continue
                    t = (shape.text or "").strip()
                    if t:
                        slide_text_count += 1
                        # Tag with slide number so search results show origin
                        parts.append(f"[SLIDE {si+1}] {t}")

            full = "\n\n".join(parts).strip()
            details["total_len"] = len(full)
            details["slides"] = len(pres.slides)
            details["text_blocks"] = slide_text_count

            return full, details
        except Exception as e:
            details["error"] = f"RUNTIME_ERROR: {type(e).__name__}: {e}"
            return "", details
